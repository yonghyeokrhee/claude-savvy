#!/usr/bin/env python3
"""Render a deck.json outline into a PPTX file."""
import argparse
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x6A, 0x73, 0x7D)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.6)).text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = INK

    sub = slide.shapes.add_textbox(Inches(0.6), Inches(3.9), Inches(12.1), Inches(1.0)).text_frame
    sub.word_wrap = True
    s = sub.paragraphs[0]
    s.text = subtitle or ""
    s.font.size = Pt(20)
    s.font.color.rgb = MUTED


def add_content_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    head = slide.shapes.add_textbox(Inches(0.95), Inches(0.45), Inches(11.8), Inches(0.9)).text_frame
    head.word_wrap = True
    p = head.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = INK

    body = slide.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(11.8), Inches(5.4)).text_frame
    body.word_wrap = True
    for i, bullet in enumerate(bullets or []):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = f"•  {bullet}"
        para.font.size = Pt(22)
        para.font.color.rgb = INK
        para.space_after = Pt(14)

    if note:
        slide.notes_slide.notes_text_frame.text = note


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--outline", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    with open(args.outline) as f:
        deck = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, deck["title"], deck.get("subtitle", ""))
    for s in deck["slides"]:
        add_content_slide(prs, s["title"], s.get("bullets", []), s.get("note"))

    prs.save(args.output)
    print(f"wrote {len(deck['slides']) + 1} slides -> {args.output}")


if __name__ == "__main__":
    main()
